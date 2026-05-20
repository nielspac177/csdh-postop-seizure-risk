"""Task 37 — Humanizer pass across manuscripts, slides and the NP working deck.

Applies a small, deliberate set of phrase-level edits to remove the most
common AI-pattern phrasings while keeping technical precision intact:

  • Slides_v2_NP.pptx          → edited in place
  • scripts/36_build_talk_slides.py  → edited in place (the build script)
  • scripts/27_build_jnnp_manuscript.py  → edited in place (build script)
  • Manuscript_05192026/Slides.pptx will be regenerated from the edited
    script 36 after this script runs

The manuscript edits are conservative: only the patterns that most strongly
read as AI-generated are touched, and only when the replacement preserves
technical meaning.  Slide edits are slightly more direct.

Run order in the pipeline:
  python scripts/37_humanizer_pass.py        # edit sources + Slides_v2_NP
  python scripts/27_build_jnnp_manuscript.py # rebuild manuscript
  python scripts/36_build_talk_slides.py     # rebuild Slides.pptx
"""
import os, sys
sys.path.insert(0, os.path.dirname(__file__))
from pathlib import Path

from pptx import Presentation

REPO = Path("/Users/nielspacheco/Desktop/Research/"
             "Ogilvy research/Data Chronic Subdural Haematoma")
NP_PPTX = REPO / "Manuscript_05192026" / "Slides_v2_NP.pptx"
SCRIPTS = REPO / "revision_analyses" / "scripts"


# ── Universal replacements applied to every text artefact ────────
# Each entry is (find, replace).  Order matters: longer / more specific
# patterns first so a broader pattern does not pre-empt them.
REPL_UNIVERSAL = [
    # AI-cliché triplets and filler
    ("It is worth noting that ", ""),
    ("it is worth noting that ", ""),
    ("Notably, ", ""),
    ("Importantly, ", ""),
    ("In this study, ", ""),
    ("In the current study, ", ""),
    ("In order to ", "To "),
    ("It should be noted that ", ""),
    ("We sought to develop and externally evaluate",
     "We developed and externally evaluated"),
    ("We sought to develop and externally validate",
     "We developed and externally validated"),
    ("We undertook a proof-of-concept study to develop and externally validate",
     "We conducted a proof-of-concept study, developing and externally validating"),
    # Slide-only punchier replacements (also safe in prose)
    ("Time to optimise something else.", "The right target is calibration."),
    ("Time to optimize something else.", "The right target is calibration."),
    ("Risk stratification is the missing piece — and is currently empirical.",
     "Risk stratification would fill this gap, and is currently empirical."),
    ("the missing piece", "the open question"),
    # Phrasing that reads as marketing
    ("regulator-friendly parametric estimators",
     "interpretable parametric estimators with valid confidence intervals"),
    ("the actionable improvement target", "the improvement target"),
    ("the actionable improvement", "the improvement"),
    ("Calibration is the actionable improvement",
     "Calibration is the improvement target"),
    ("no over-claim", "without over-reaching"),
    # Triplet over-enumeration in slides
    ("Five conclusions in one minute", "Five take-homes"),
    # Hedge softening
    ("To our knowledge this is the first",
     "This is, to our knowledge, the first"),
    ("To our knowledge, this is the first",
     "This is, to our knowledge, the first"),
    # Subtle but recurring AI tropes
    ("We applied conformal prediction sets",
     "Conformal prediction sets were applied"),
    ("supplying the deployment with a principled abstention signal",
     "providing a principled abstention signal at the bedside"),
    # Idiosyncratic but cleaner
    ("the deployment with a principled abstention signal",
     "a principled abstention signal at the bedside"),
    ("decision-relevant value, not by p-value",
     "their decision-relevant value rather than by p-value alone"),
    # Casual contractions in the slides → formal
    ("I don't know", "I do not know"),
    ("don't know", "do not know"),
    # Stylistic — slightly over-formal "honestly deployable" reads as AI
    ("can be honestly deployable when",
     "can be deployable when"),
    ("honestly deployable", "deployable"),
    # Buzz / cliché softening
    ("stress-test", "test"),
    # Sentence-opener variety: avoid stacking "We ..." starts (handled per-script below)
]


def humanize_text(s):
    for find, repl in REPL_UNIVERSAL:
        s = s.replace(find, repl)
    return s


def humanize_pptx_in_place(path):
    """Open a .pptx, run the humanizer on every text run, save back."""
    prs = Presentation(path)
    n_edits = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame: continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    new = humanize_text(run.text)
                    if new != run.text:
                        run.text = new
                        n_edits += 1
    prs.save(path)
    return n_edits


def humanize_script_in_place(path):
    """Open a .py source file, apply REPL_UNIVERSAL to its string contents,
    and write back.  Edits are textual at the source level — they will
    propagate to every artefact built from the script."""
    src = Path(path).read_text(encoding="utf-8")
    new = src
    edits = 0
    for find, repl in REPL_UNIVERSAL:
        before = new
        new = new.replace(find, repl)
        if new != before:
            edits += before.count(find)
    if new != src:
        Path(path).write_text(new, encoding="utf-8")
    return edits


def main():
    print("Humanizer pass — phrase-level cleanup")
    print()
    print(f"[1/3] Editing manuscript build script (script 27)")
    n27 = humanize_script_in_place(SCRIPTS / "27_build_jnnp_manuscript.py")
    print(f"      {n27} edits applied")
    print(f"[2/3] Editing slide build script (script 36)")
    n36 = humanize_script_in_place(SCRIPTS / "36_build_talk_slides.py")
    print(f"      {n36} edits applied")
    print(f"[3/3] Editing Slides_v2_NP.pptx in place")
    if NP_PPTX.exists():
        nNP = humanize_pptx_in_place(NP_PPTX)
        print(f"      {nNP} run-level text replacements")
    else:
        print(f"      [SKIP] {NP_PPTX} not found")
    print()
    print("Done.  Rebuild artefacts:")
    print("  python scripts/27_build_jnnp_manuscript.py")
    print("  python scripts/36_build_talk_slides.py")
    print("  python scripts/32_build_code_companion.py")
    print("  python scripts/33_build_code_companion_pdf.py")


if __name__ == "__main__":
    main()
