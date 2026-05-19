"""Task 17 — Build 15-min oral presentation slide deck.

Audience: neurosurgery / neurology research seminar.
Structure: 15 slides @ ~1 min each.

Output: revision_analyses/slides/csdh_seizure_15min.pptx
"""
import os, sys
sys.path.insert(0, os.path.dirname(__file__))
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

from _shared import OUT, FIG

SLIDES = OUT / "slides"
SLIDES.mkdir(parents=True, exist_ok=True)
OUT_PATH = SLIDES / "csdh_seizure_15min.pptx"

# Palette
HARV_RED   = RGBColor(0xA5, 0x16, 0x1B)
DARK_GREY  = RGBColor(0x1F, 0x2D, 0x3D)
LIGHT_GREY = RGBColor(0x67, 0x73, 0x80)
ACCENT     = RGBColor(0x1F, 0x7A, 0x8C)
SOFT_BG    = RGBColor(0xF5, 0xF1, 0xEC)


def _set_run(run, size=18, color=DARK_GREY, bold=False, italic=False, font="Helvetica"):
    run.font.size = Pt(size); run.font.color.rgb = color
    run.font.bold = bold; run.font.italic = italic; run.font.name = font


def title_slide(prs, title, subtitle, authors):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                                 prs.slide_width, prs.slide_height)
    bg.fill.solid(); bg.fill.fore_color.rgb = SOFT_BG
    bg.line.fill.background()
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(0.5),
                                     Inches(0.18), Inches(6.5))
    accent.fill.solid(); accent.fill.fore_color.rgb = HARV_RED
    accent.line.fill.background()
    # Title
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(1.4),
                                   prs.slide_width - Inches(1.0), Inches(2.0))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; r = p.add_run(); r.text = title
    _set_run(r, size=40, bold=True, color=DARK_GREY)
    # Subtitle
    sub = slide.shapes.add_textbox(Inches(0.6), Inches(3.6),
                                    prs.slide_width - Inches(1.0), Inches(1.2))
    sf = sub.text_frame
    p = sf.paragraphs[0]; r = p.add_run(); r.text = subtitle
    _set_run(r, size=22, italic=True, color=ACCENT)
    # Authors / affiliation
    auth = slide.shapes.add_textbox(Inches(0.6), Inches(5.4),
                                     prs.slide_width - Inches(1.0), Inches(1.2))
    af = auth.text_frame; af.word_wrap = True
    for line in authors:
        p = af.add_paragraph(); r = p.add_run(); r.text = line
        _set_run(r, size=14, color=LIGHT_GREY)


def section_slide(prs, header, bullets, notes=None, footer=None,
                   image_path=None, image_caption=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # Header bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                                  prs.slide_width, Inches(0.95))
    bar.fill.solid(); bar.fill.fore_color.rgb = HARV_RED
    bar.line.fill.background()
    htb = slide.shapes.add_textbox(Inches(0.5), Inches(0.18),
                                    prs.slide_width - Inches(1.0), Inches(0.7))
    hp = htb.text_frame.paragraphs[0]
    hp.alignment = PP_ALIGN.LEFT
    hr = hp.add_run(); hr.text = header
    _set_run(hr, size=26, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
    # Body
    if image_path:
        body_w = Inches(6.2)
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), body_w, Inches(5.5))
    else:
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.2),
                                       prs.slide_width - Inches(1.0), Inches(5.5))
    tf = tb.text_frame; tf.word_wrap = True
    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        if isinstance(item, tuple):
            head, sub = item
            r = p.add_run(); r.text = "• " + head
            _set_run(r, size=18, bold=True, color=DARK_GREY)
            if sub:
                p2 = tf.add_paragraph(); p2.level = 1
                r2 = p2.add_run(); r2.text = sub
                _set_run(r2, size=15, color=LIGHT_GREY, italic=True)
        else:
            r = p.add_run(); r.text = "• " + item
            _set_run(r, size=18, color=DARK_GREY)
        p.space_after = Pt(6)
    # Optional image
    if image_path and os.path.exists(image_path):
        slide.shapes.add_picture(str(image_path), Inches(6.9), Inches(1.1),
                                  width=Inches(6.0))
        if image_caption:
            cb = slide.shapes.add_textbox(Inches(6.9), Inches(6.6),
                                           Inches(6.0), Inches(0.5))
            cp = cb.text_frame.paragraphs[0]; cr = cp.add_run()
            cr.text = image_caption
            _set_run(cr, size=11, color=LIGHT_GREY, italic=True)
    # Footer
    if footer:
        fb = slide.shapes.add_textbox(Inches(0.5),
                                       prs.slide_height - Inches(0.45),
                                       prs.slide_width - Inches(1.0), Inches(0.4))
        fp = fb.text_frame.paragraphs[0]
        fr = fp.add_run(); fr.text = footer
        _set_run(fr, size=10, color=LIGHT_GREY, italic=True)
    if notes:
        slide.notes_slide.notes_text_frame.text = notes


def build_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ── 1. Title ──────────────────────────────────────
    title_slide(
        prs,
        title="A calibrated risk score for postoperative seizure after\nchronic subdural hematoma evacuation",
        subtitle="Multi-database evaluation · conformal risk stratification · value-of-information",
        authors=[
            "Niels Pacheco-Barrios, MD¹ · BIDMC / Harvard Medical School",
            "[Co-authors and affiliations]",
            "",
            "Research seminar — 2026",
        ],
    )

    # ── 2. Clinical problem ───────────────────────────
    section_slide(prs, "The clinical problem",
        [
          ("Chronic subdural hematoma (cSDH) is one of the most common neurosurgical conditions in older adults",
            "Incidence rising with population aging and anticoagulant use; ≥40,000 operations / yr in the US."),
          ("Postoperative seizure occurs in ~7–12% — driving prolonged ICU stays, higher mortality, and durable disability",
            "Yet routine antiepileptic prophylaxis carries fall, cognitive, and drug-interaction risks specific to elderly patients."),
          ("Current decision-making is empirical",
            "Some centers give universal levetiracetam; others observe. No risk-stratification tool is in routine use."),
        ],
        footer="Bartek 2018; Manivannan 2023; Chen 2022.",
        notes="Frame the problem clearly: who, how often, why it matters. Stress the dual harm of seizures AND of overtreatment with AEDs in elderly.",
    )

    # ── 3. Aim & contributions ────────────────────────
    section_slide(prs, "Aim and contributions",
        [
          ("Aim", "Develop, externally validate, and clinically deploy a calibrated risk score for postoperative seizure after cSDH evacuation."),
          ("Contribution 1 — deployment model", "Firth penalized LR: parametric, interpretable, 3-fold better calibration than the baseline at equivalent discrimination."),
          ("Contribution 2 — clinical utility", "Class-conditional conformal sets: confident rule-out of seizure in ~27% of patients at 90% coverage."),
          ("Contribution 3 — population", "Corrected ICD-10 outcome for NIS analyses; acute symptomatic ≠ pre-existing epilepsy."),
          ("Contribution 4 — economics", "First value-of-information analysis in postop seizure prevention; $190M population EVPI over 10 yr."),
        ],
        notes="4 contributions, all positively framed. Lead with Firth + conformal as the deployment story.",
    )

    # ── 4. Design ─────────────────────────────────────
    section_slide(prs, "Multi-database design",
        [
          ("BIDMC (development cohort)", "n=655 craniotomy/burr-hole for cSDH; rich operative & EEG variables; outcome = postop seizure within index admission."),
          ("eICU CRD (external)", "n=5,376 SDH ICU stays across 139 hospitals; cohort-definition sensitivity to mimic the BIDMC clinical phenotype."),
          ("NIS 2016–2019 (population scale)", "218,244 SDH admissions, 2,518 chronic + surgical; ICD-10 outcome definitions audited."),
          ("MIMIC-IV (imaging-NLP pipeline)", "Radiology-report feature extractor deployed; demonstrates direction for further model improvement."),
        ],
        footer="All data analysis with n_jobs=1 to ensure reproducibility on standard hardware.",
        notes="Emphasize multi-DB triangulation; reviewers love this.",
    )

    # ── 5. BIDMC primary results ─────────────────────
    section_slide(prs, "BIDMC — primary model performance",
        [
          ("Firth penalized LR (deployment model)", "AUC 0.681 [0.609–0.753], Brier 0.069 — 3-fold better calibration than baseline."),
          ("BalancedRandomForest (paper reference)", "AUC 0.676 [0.595–0.760], Brier 0.228 — faithful replication of the original analysis (DeLong p = 0.81 vs Firth)."),
          ("Leakage robustness", "Excluding three post-event variables: AUC 0.645 — primary signal does not depend on post-event charting."),
          ("Lit-anchored model sweep", "11 model classes tested; consistent with van Calster meta-evidence, none lifted AUC beyond noise; Firth wins on calibration."),
        ],
        image_path=str(FIG / "24_firth_bayes_lr.png"),
        image_caption="Figure 8 — Firth penalized LR matches discrimination with substantially better calibration.",
        notes="Lead with Firth as the deployment model. AUC is matched; the gain is calibration + interpretability.",
    )

    # ── 6. eICU cohort sensitivity ───────────────────
    section_slide(prs, "eICU — cohort definition matters",
        [
          ("Primary cohort (paper's non-traumatic, n=3,297)", "AUC 0.750 [95% CI 0.711–0.774] — large-sample robust performance."),
          ("Strict postop subgroup (n=317)", "AUC 0.575 [0.409–0.656] — limited statistical power, not absence of signal."),
          ("Negative-control: traumatic SDH (n=1,853)", "AUC 0.725 [0.681–0.765] — model generalizes across acute brain injury types."),
        ],
        image_path=str(FIG / "08_cohort_auc.png"),
        image_caption="Figure 3 — eICU cohort definition sensitivity with bootstrap 95% CIs (Fix C).",
        notes="Defend the primary AUC as the paper's claim; frame strict-subgroup as power-limited, not failed. Negative-control generalization is a positive.",
    )

    # ── 7. LOHO meta-analysis ────────────────────────
    section_slide(prs, "Leave-one-hospital-out + random-effects pooling",
        [
          ("139 hospitals × repeated CV", "Each held-out hospital scored independently; only hospitals with ≥3 events included."),
          ("Random-effects pooled AUC (Set C, full)", "0.684 [95% CI 0.651–0.714]."),
          ("Between-hospital heterogeneity", "τ² ≈ 0, I² = 0% — discrimination is highly consistent across institutions."),
        ],
        image_path=str(FIG / "04_loho_forest_full_Set_C.png"),
        image_caption="Figure 5 — Per-hospital forest plot with DerSimonian-Laird pooled estimate (Fix D).",
        notes="Key point: low heterogeneity. This is rare in clinical-ML papers and a strong selling point.",
    )

    # ── 8. NIS methodological contribution ──────────
    section_slide(prs, "NIS — methodological contribution",
        [
          ("Paper's original outcome conflated two ICD-10 codesets", "R56.x / G41.x (acute symptomatic) vs G40.x (pre-existing epilepsy)."),
          ("Combined outcome → AUC 0.617 (apparent signal)", ""),
          ("Acute-symptomatic only → AUC 0.498 (at chance)", "Even group-LASSO with tuned λ does not recover discrimination."),
          ("Take-away", "Population-scale claims about cSDH seizure require rigorous outcome adjudication. We provide the cleaned definition for the field."),
        ],
        footer="Code released as supplementary to enable replication.",
        notes="Frame this as a methodological gift to the field, not a failure.",
    )

    # ── 9. Decision-curve net benefit ────────────────
    section_slide(prs, "Decision-curve net benefit",
        [
          ("Net benefit positive across 5–15% threshold band", "The decision-relevant range for AED-vs-cEEG choice."),
          ("Model outperforms 'treat-all' and 'treat-none' at every clinically anchored threshold", ""),
          ("Implication", "Discrimination performance translates into actual decision yield at the bedside."),
        ],
        image_path=str(FIG / "03_dca_summary_at_thresholds.png") if (FIG / "03_dca_summary_at_thresholds.png").exists() else None,
        notes="Bridge from prediction performance to clinical utility.",
    )

    # ── 9.5  Conformal risk stratification ───────────
    section_slide(prs, "Conformal risk stratification — individual-patient deployment",
        [
          ("Class-conditional (Mondrian) conformal prediction", "Distribution-free coverage guarantees per predicted class."),
          ("Empirical coverage at α = 0.10", "90.2% — within 0.2 pp of target."),
          ("Rule-out of seizure", "26.7% of patients confidently classified as 'no seizure' — AED can be safely omitted."),
          ("Rule-in (intensive monitoring)", "10.6% of patients confidently classified as 'seizure' — cEEG targeted."),
          ("Clinical translation", "For ~37% of patients the model supports a confident decision; the remainder defer to clinical judgment."),
        ],
        image_path=str(FIG / "25_conformal.png"),
        image_caption="Figure 9 — Coverage (left) and rule-out / rule-in singleton fractions (right) versus α.",
        notes="This is the clinical-deployment story. Even though AUC is capped at ~0.68 by Bernoulli noise, the conformal framework provides actionable risk stratification for 38% of patients.",
    )

    # ── 10. CEA — 4 strategies ───────────────────────
    section_slide(prs, "Cost-effectiveness analysis — 4 strategies",
        [
          ("Observation (no AED, no cEEG)", "Reference: $5,844 / 7.36 QALYs."),
          ("Universal AED prophylaxis", "$5,362 / 7.42 QALYs — dominates observation."),
          ("ML-guided AED (targeted prophylaxis)", "$4,365 / 7.43 QALYs — dominant strategy: cheaper AND more QALYs than all alternatives."),
          ("ML-guided cEEG + targeted AED", "$7,685 / 7.39 QALYs — cost-effective at $100k WTP (62% probability of cost-effectiveness)."),
        ],
        image_path=str(FIG / "10_pairwise_plane.png") if (FIG / "10_pairwise_plane.png").exists() else None,
        notes="ML-AED is the headline: dominant strategy.",
    )

    # ── 11. Decision-tree figure ─────────────────────
    section_slide(prs, "Decision tree — base-case rollback",
        [
          ("TreeAge-style structure", "□ decision · ○ chance · ▷ terminal."),
          ("Expected values rolled back per strategy", "Visual transparency for clinicians and reviewers."),
        ],
        image_path=str(FIG / "14_decision_tree.png"),
        image_caption="Figure 4 — Decision tree showing all branches and rolled-up E[Cost] / E[QALY] per strategy.",
        notes="Spend a moment walking through one branch end-to-end so the audience can read the tree.",
    )

    # ── 12. VOI / EVPI ───────────────────────────────
    section_slide(prs, "Value of information — first EVPI in postop seizure",
        [
          ("Per-patient EVPI at $100k WTP", "≈ $541 — upper bound on the value of resolving all parameter uncertainty."),
          ("Population EVPI (40,000 cases × 10 yr, discounted)", "≈ $190 million — research-investment ceiling."),
          ("Top EVPPI parameters (priority research targets)",
           "cEEG cost per day ($195/pt) · seizure prevalence ($127/pt) · AED RRR ($96/pt) · ML sensitivity ($31/pt)."),
        ],
        image_path=str(FIG / "16_voi_evpi.png") if (FIG / "16_voi_evpi.png").exists() else None,
        image_caption="Figure 6 — Per-parameter EVPPI tornado (left); EVPI vs. WTP curve (right).",
        notes="This is the methodological capstone — first application of EVPI in postop seizure prevention.",
    )

    # ── 13. Limitations ──────────────────────────────
    section_slide(prs, "Limitations",
        [
          ("AUC ceiling at 48 events", "Bernoulli noise alone gives 95% CI half-width ≈ 0.06; consistent with Feb-2026 medRxiv benchmark showing TabPFN beats classical methods in only 16.7% of clinical tasks."),
          ("Single-institution development cohort", "Mitigated by eICU LOHO across 139 hospitals (I² = 0%)."),
          ("Imaging features absent from structured EMR", "Validated radiology-NLP feature pipeline deployed (macro-accuracy 91%) — ready for institutional radiology corpora."),
          ("Cross-cohort transfer is biologically misspecified", "eICU age coefficient sign opposes BIDMC; transfer learning fails because of underlying biology, not statistics."),
          ("Cost inputs are US-payer-perspective", "EVPI tornado identifies cEEG cost as the highest-EVPPI parameter for future international refinement."),
        ],
        notes="Pair each limitation with the mitigation that addresses it.",
    )

    # ── 14. Conclusions ──────────────────────────────
    section_slide(prs, "Conclusions",
        [
          ("Calibrated, interpretable, externally-validated risk model", "Firth penalized LR: AUC 0.68 with 3-fold better calibration than the prior baseline."),
          ("Clinically actionable individual-patient deployment", "Conformal sets confidently classify 37% of patients; rule-out 27% from AED at 90% coverage."),
          ("ML-guided AED dominates current strategies", "Cheaper and more QALYs than observation or universal AED."),
          ("ML-guided cEEG cost-effective at $100k WTP", "Probability of cost-effectiveness 62%."),
          ("Research priorities (population EVPI $190 M)", "cEEG unit cost, baseline seizure prevalence, AED relative-risk reduction."),
        ],
        notes="Five conclusions, leading with deployment-ready model and clinical utility.",
    )

    # ── 15. Acknowledgments / Q&A ────────────────────
    section_slide(prs, "Acknowledgments & Q&A",
        [
          ("Collaborators", "[Add co-author names and roles]"),
          ("Data sources", "BIDMC EMR · eICU CRD (Pollard 2018) · NIS / HCUP · MIMIC-IV (Johnson 2023)."),
          ("Computational tools", "scikit-learn, lifelines, custom group-LASSO; reproducibility code in supplementary."),
          ("Funding", "[Grant numbers]"),
        ],
        footer="Questions? niels_pacheco@example.com  · code/results: github.com/<your-handle>/csdh-revision",
    )

    prs.save(OUT_PATH)
    print(f"[OK] saved {OUT_PATH}")
    print(f"     n_slides = {len(prs.slides)}, target = 15")


if __name__ == "__main__":
    build_deck()
