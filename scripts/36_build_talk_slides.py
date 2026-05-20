"""Task 36 — 15-minute talk slide deck.

Story-arc structure that follows IMRAD but with a deliberate narrative
beat at the centre:

  1. Title (with QR code to GitHub repo)
  2-3.  Introduction — clinical problem + decision controversy + ML gap
  4-5.  Methods — multi-database design + modelling strategy
  6.    Results turn 1 — the AUC ceiling is biological, not modelling failure
  7.    Results turn 2 — calibration is the actionable target
  8.    Results turn 3 — conformal prediction translates into bedside decisions
  9.    Results turn 4 — ML-AED dominates current practice (CEA)
  10.   Results turn 5 — value-of-information identifies research priorities
  11.   Discussion — methodological contributions
  12.   Limitations — honest acknowledgement, paired with mitigations
  13.   Conclusions — the five take-homes
  14.   Acknowledgements + QR

Output: Manuscript_05192026/Slides.pptx
"""
import os, sys
sys.path.insert(0, os.path.dirname(__file__))
from pathlib import Path

import qrcode
from PIL import Image

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

from _shared import FIG

SUB_DIR = Path("/Users/nielspacheco/Desktop/Research/"
                 "Ogilvy research/Data Chronic Subdural Haematoma/"
                 "Manuscript_05192026")
SUB_DIR.mkdir(parents=True, exist_ok=True)
OUT_PATH = SUB_DIR / "Slides.pptx"

REPO_URL = "https://github.com/nielspac177/csdh-postop-seizure-risk"
SITE_URL = "https://nielspac177.github.io/csdh-postop-seizure-risk"

# ── Palette (matches manuscript figures) ───────────────────────────
NAVY      = RGBColor(0x1F, 0x3D, 0x5C)
NAVY_DARK = RGBColor(0x0F, 0x24, 0x3A)
NAVY_SOFT = RGBColor(0xEE, 0xF3, 0xF8)
RUST      = RGBColor(0xB5, 0x53, 0x2C)
FOREST    = RGBColor(0x2E, 0x6B, 0x45)
OCHRE     = RGBColor(0xB5, 0x8A, 0x2E)
INK       = RGBColor(0x26, 0x23, 0x20)
INK_50    = RGBColor(0xF5, 0xF3, 0xEE)
GREY      = RGBColor(0x6E, 0x6E, 0x6E)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)

# ── QR code (generated once, embedded in title + final slide) ─────
QR_PATH = SUB_DIR / "_slides_qr.png"

def make_qr(url, out, size=600, fg="#1F3D5C", bg="#FFFFFF"):
    qr = qrcode.QRCode(version=4, error_correction=qrcode.constants.ERROR_CORRECT_M,
                        box_size=12, border=2)
    qr.add_data(url); qr.make(fit=True)
    img = qr.make_image(fill_color=fg, back_color=bg).convert("RGB")
    if img.size[0] != size: img = img.resize((size, size), Image.LANCZOS)
    img.save(out)
    return out

make_qr(REPO_URL, QR_PATH)


# ── Style helpers ─────────────────────────────────────────────────
def _set_run(run, *, size=18, color=INK, bold=False, italic=False,
              font="Helvetica"):
    run.font.size = Pt(size); run.font.name = font
    run.font.color.rgb = color
    run.font.bold = bold; run.font.italic = italic


def blank_slide(prs, *, bg=INK_50):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg_rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                                       prs.slide_width, prs.slide_height)
    bg_rect.fill.solid(); bg_rect.fill.fore_color.rgb = bg
    bg_rect.line.fill.background()
    return slide


def add_box(slide, x, y, w, h, *, fill=None, border=None, border_w=1.0):
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    if fill is None:
        box.fill.background()
    else:
        box.fill.solid(); box.fill.fore_color.rgb = fill
    if border is None:
        box.line.fill.background()
    else:
        box.line.color.rgb = border
        box.line.width = Pt(border_w)
    return box


def add_text(slide, text, x, y, w, h, *, size=18, color=INK, bold=False,
              italic=False, font="Helvetica", align=PP_ALIGN.LEFT,
              fill=None, line_spacing=1.15):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.03)
    if fill is not None:
        tb.fill.solid(); tb.fill.fore_color.rgb = fill
        tb.line.fill.background()
    p = tf.paragraphs[0]
    p.alignment = align; p.line_spacing = line_spacing
    r = p.add_run(); r.text = text
    _set_run(r, size=size, color=color, bold=bold, italic=italic, font=font)
    return tb


def add_header(slide, prs, eyebrow, title, *, accent=NAVY):
    # Top stripe
    add_box(slide, 0, 0, prs.slide_width, Inches(0.42), fill=accent)
    add_text(slide, eyebrow.upper(), Inches(0.5), Inches(0.06),
              prs.slide_width - Inches(1.0), Inches(0.32),
              size=11, color=WHITE, bold=True, font="Helvetica")
    # Title
    add_text(slide, title, Inches(0.5), Inches(0.55),
              prs.slide_width - Inches(1.0), Inches(0.85),
              size=28, color=NAVY_DARK, bold=True, font="Helvetica")
    # Subtle underline
    add_box(slide, Inches(0.5), Inches(1.30), Inches(2.0), Inches(0.04),
              fill=accent)


def add_footer(slide, prs, *, slide_num=None, total=14):
    add_text(slide, "Pacheco-Barrios  ·  csdh-postop-seizure-risk",
              Inches(0.5), prs.slide_height - Inches(0.40),
              Inches(7), Inches(0.30),
              size=9, color=GREY, font="Helvetica")
    if slide_num is not None:
        add_text(slide, f"{slide_num} / {total}",
                  prs.slide_width - Inches(1.2),
                  prs.slide_height - Inches(0.40),
                  Inches(1.0), Inches(0.30),
                  size=9, color=GREY, align=PP_ALIGN.RIGHT, font="Helvetica")


def add_bullets(slide, items, x, y, w, h, *, size=18, color=INK,
                  bullet_size=18, line_spacing=1.4):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        if isinstance(item, tuple):
            head, sub = item
        else:
            head, sub = item, None
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        p.space_after = Pt(8)
        # bullet diamond
        rd = p.add_run(); rd.text = "◆  "
        _set_run(rd, size=bullet_size, color=NAVY, bold=True)
        rh = p.add_run(); rh.text = head
        _set_run(rh, size=size, color=color, bold=True)
        if sub:
            p2 = tf.add_paragraph()
            p2.alignment = PP_ALIGN.LEFT
            p2.line_spacing = line_spacing
            p2.space_after = Pt(8)
            p2.level = 1
            rs = p2.add_run(); rs.text = sub
            _set_run(rs, size=size - 2, color=GREY)
    return tb


def add_image(slide, path, x, y, w=None, h=None):
    if not os.path.exists(path): return None
    return slide.shapes.add_picture(str(path), x, y, width=w, height=h)


def add_qr(slide, x, y, side, *, caption=None):
    pic = slide.shapes.add_picture(str(QR_PATH), x, y, width=side, height=side)
    if caption:
        add_text(slide, caption,
                  x - Inches(0.10), y + side + Inches(0.04),
                  side + Inches(0.20), Inches(0.40),
                  size=10, color=GREY, italic=True, align=PP_ALIGN.CENTER)
    return pic


def add_callout(slide, text, x, y, w, h, *, accent=RUST):
    box = add_box(slide, x, y, w, h, fill=INK_50, border=accent, border_w=1.4)
    tb = slide.shapes.add_textbox(x + Inches(0.20), y + Inches(0.15),
                                    w - Inches(0.40), h - Inches(0.30))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT; p.line_spacing = 1.25
    r = p.add_run(); r.text = text
    _set_run(r, size=14, color=INK, italic=True)
    return box


# ── Slide builders ────────────────────────────────────────────────
def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ─────── 1. Title slide ───────
    s = blank_slide(prs, bg=INK_50)
    # Accent ribbon down the left side
    add_box(s, 0, 0, Inches(0.30), prs.slide_height, fill=NAVY)
    # Eyebrow tag
    add_text(s, "PROOF-OF-CONCEPT  ·  RESEARCH SEMINAR",
              Inches(0.7), Inches(0.6), Inches(10), Inches(0.4),
              size=12, color=RUST, bold=True)
    # Main title
    add_text(s,
              "A calibrated risk score for postoperative seizure\n"
              "after chronic subdural haematoma evacuation",
              Inches(0.7), Inches(1.1), Inches(10), Inches(2.1),
              size=36, color=NAVY_DARK, bold=True, line_spacing=1.18)
    # Subtitle
    add_text(s,
              "A proof-of-concept multi-database study with value-of-information analysis",
              Inches(0.7), Inches(3.4), Inches(11), Inches(0.7),
              size=18, color=GREY, italic=True)
    # Author
    add_text(s, "Niels Pacheco-Barrios, MD",
              Inches(0.7), Inches(4.6), Inches(8), Inches(0.5),
              size=18, color=INK, bold=True)
    add_text(s,
              "Department of Neurosurgery  ·  Beth Israel Deaconess Medical Center\n"
              "Harvard Medical School",
              Inches(0.7), Inches(5.05), Inches(8), Inches(1.0),
              size=13, color=GREY, italic=True, line_spacing=1.4)
    # Bottom strip
    add_box(s, 0, prs.slide_height - Inches(0.60), prs.slide_width, Inches(0.60),
              fill=NAVY)
    add_text(s, "Code, data-availability and interactive companion site →",
              Inches(0.7), prs.slide_height - Inches(0.50),
              Inches(8), Inches(0.40),
              size=12, color=WHITE, bold=True)
    add_text(s, "github.com/nielspac177/csdh-postop-seizure-risk",
              Inches(0.7), prs.slide_height - Inches(0.28),
              Inches(8), Inches(0.30),
              size=10, color=NAVY_SOFT, italic=True)
    # QR
    add_qr(s, prs.slide_width - Inches(1.85), Inches(4.6),
            Inches(1.55))
    add_text(s, "Scan to explore",
              prs.slide_width - Inches(1.95), Inches(6.25),
              Inches(1.75), Inches(0.30),
              size=10, color=GREY, italic=True, align=PP_ALIGN.CENTER)

    # ─────── 2. Why this matters (introduction 1/3) ───────
    s = blank_slide(prs)
    add_header(s, prs, "Introduction · 1/3", "Why this matters")
    add_bullets(s, [
        ("Chronic subdural haematoma is now one of the most common neurosurgical conditions in older adults",
         "Annual operative incidence in the United States approaches 40,000 cases and is rising with population ageing and anticoagulant use."),
        ("Postoperative seizure complicates 7–12% of cSDH evacuations",
         "Independently associated with prolonged ICU stay, 30-day mortality, and lower functional independence at follow-up."),
        ("Yet routine antiepileptic prophylaxis carries specific harm in the elderly",
         "Levetiracetam: neuropsychiatric AE ~15–20%, somnolence ~28%, fall RR 1.6–1.8 vs no AED. The very population most likely to need a seizure-prevention strategy is also the most vulnerable to its side effects."),
    ], Inches(0.6), Inches(1.6), Inches(12.1), Inches(5.3))
    add_footer(s, prs, slide_num=2)

    # ─────── 3. The decision dilemma (introduction 2/3) ───────
    s = blank_slide(prs)
    add_header(s, prs, "Introduction · 2/3", "The clinical decision lives between two harms")
    # Two columns: AED vs cEEG
    col_w = Inches(5.7); col_h = Inches(4.8); col_y = Inches(1.7)
    # Left column - AED
    add_box(s, Inches(0.6), col_y, col_w, col_h, fill=WHITE,
              border=NAVY, border_w=1.5)
    add_text(s, "AED prophylaxis for all",
              Inches(0.85), col_y + Inches(0.25), col_w - Inches(0.5), Inches(0.5),
              size=18, color=NAVY, bold=True)
    add_text(s,
              "Empirical · cheap · simple\n\n"
              "But in the elderly:\n"
              "  •  Falls (RR 1.6–1.8)\n"
              "  •  Cognitive impairment, somnolence\n"
              "  •  Drug interactions, polypharmacy\n"
              "  •  Often continued beyond hospital stay",
              Inches(0.85), col_y + Inches(0.85),
              col_w - Inches(0.5), col_h - Inches(1.1),
              size=14, color=INK, line_spacing=1.5)
    # Right column - cEEG
    add_box(s, Inches(7.0), col_y, col_w, col_h, fill=WHITE,
              border=RUST, border_w=1.5)
    add_text(s, "Selective cEEG monitoring",
              Inches(7.25), col_y + Inches(0.25), col_w - Inches(0.5), Inches(0.5),
              size=18, color=RUST, bold=True)
    add_text(s,
              "Avoids unnecessary AED · captures subclinical seizures\n\n"
              "But:\n"
              "  •  Capacity-limited in many centres\n"
              "  •  ~$1,500–3,000 per monitoring episode\n"
              "  •  Requires interpretation expertise\n"
              "  •  Who do we select?",
              Inches(7.25), col_y + Inches(0.85),
              col_w - Inches(0.5), col_h - Inches(1.1),
              size=14, color=INK, line_spacing=1.5)
    add_callout(s,
                 "Risk stratification is the missing piece — and is currently empirical.",
                 Inches(0.6), Inches(6.65), Inches(12.1), Inches(0.55),
                 accent=OCHRE)
    add_footer(s, prs, slide_num=3)

    # ─────── 4. The gap prior ML hasn't closed (introduction 3/3) ───────
    s = blank_slide(prs)
    add_header(s, prs, "Introduction · 3/3",
                "What prior ML attempts have not delivered")
    add_bullets(s, [
        ("Small single-institution cohorts",
         "Many cSDH seizure-prediction models in the literature report on n<300 with no external test."),
        ("Discrimination only, no calibration",
         "AUC alone does not say whether a probability is trustworthy at the bedside."),
        ("No clinical translation layer",
         "A probability is not a decision: how confident are we in this patient specifically?"),
        ("No decision-analytic integration",
         "Is the model worth deploying? Cost-effective at what willingness-to-pay? Which parameters drive uncertainty?"),
    ], Inches(0.6), Inches(1.6), Inches(12.1), Inches(4.6),
       size=17)
    add_callout(s,
                 "Our aim: a proof-of-concept that closes all four gaps simultaneously, "
                 "by replacing AUC with calibration + decision integration as the optimisation target.",
                 Inches(0.6), Inches(6.45), Inches(12.1), Inches(0.75),
                 accent=NAVY)
    add_footer(s, prs, slide_num=4)

    # ─────── 5. Multi-database design (methods 1/2) ───────
    s = blank_slide(prs)
    add_header(s, prs, "Methods · 1/2", "Multi-database design")
    add_image(s, FIG / "F0_graphical_abstract.png",
                Inches(0.45), Inches(1.7), w=Inches(12.4))
    add_callout(s,
                 "BIDMC develops the deployment model; eICU's 139 hospitals stress-test "
                 "external generalisability; NIS surfaces a methodological correction "
                 "needed for population-scale analyses.",
                 Inches(0.6), Inches(6.50), Inches(12.1), Inches(0.75),
                 accent=NAVY)
    add_footer(s, prs, slide_num=5)

    # ─────── 6. Modelling strategy (methods 2/2) ───────
    s = blank_slide(prs)
    add_header(s, prs, "Methods · 2/2", "An eleven-method modelling battery")
    add_bullets(s, [
        ("Baseline reference: BalancedRandomForest (the prior baseline)",
         "Reproduces the original analysis exactly so every comparison is on a like-for-like footing."),
        ("Deployment model: Firth penalized logistic regression",
         "Bias-corrected for rare events (Firth 1993, Puhr 2017). Parametric, interpretable, valid coefficient CIs."),
        ("Sensitivity battery: six SMOTE-family oversamplers, two Bayesian variants, Optuna-tuned XGBoost / LightGBM, diverse-base stacking",
         "All evaluated by paired DeLong tests against the baseline; all 5×5 repeated stratified CV."),
        ("Clinical-utility layer: class-conditional (Mondrian) conformal prediction",
         "Vovk 2005; Angelopoulos & Bates 2021. Distribution-free coverage guarantees per predicted class."),
        ("Decision-analytic layer: 4-strategy decision tree → 10-year Markov → 10,000-iteration PSA → EVPI / EVPPI",
         "Strong–Oakley 2014 non-parametric regression on 5,000 PSA samples for per-parameter information value."),
    ], Inches(0.6), Inches(1.6), Inches(12.1), Inches(5.5), size=14)
    add_footer(s, prs, slide_num=6)

    # ─────── 7. The discrimination ceiling is biological (results 1/5) ───────
    s = blank_slide(prs)
    add_header(s, prs, "Results · 1/5", "The discrimination ceiling is biological")
    add_image(s, FIG / "F3_method_battery.png",
                Inches(0.5), Inches(1.6), w=Inches(8.2))
    add_text(s,
              "Eleven methods converge at AUC ≈ 0.68",
              Inches(9.0), Inches(1.7), Inches(4.1), Inches(0.5),
              size=18, color=NAVY, bold=True)
    add_bullets(s, [
        ("48 events → 95% CI half-width on AUC=0.70 is ≈0.06 from Bernoulli noise alone", None),
        ("No method statistically lifts AUC above baseline (DeLong p>0.05 everywhere)", None),
        ("Consistent with 2022–2025 meta-evidence (van den Goorbergh, Carriero, Piccininni) that class-imbalance corrections do not raise AUC in clinical risk models", None),
    ], Inches(9.0), Inches(2.3), Inches(4.0), Inches(4.6),
       size=12, bullet_size=12)
    add_callout(s,
                 "The ceiling is biological, not modelling failure. Time to optimise something else.",
                 Inches(0.5), Inches(6.50), Inches(12.3), Inches(0.75),
                 accent=RUST)
    add_footer(s, prs, slide_num=7)

    # ─────── 8. Calibration is the actionable target (results 2/5) ───────
    s = blank_slide(prs)
    add_header(s, prs, "Results · 2/5", "Calibration is the actionable improvement")
    add_image(s, FIG / "F2_calibration_dca.png",
                Inches(0.5), Inches(1.6), w=Inches(8.2))
    add_text(s, "Firth penalized LR",
              Inches(9.0), Inches(1.7), Inches(4.1), Inches(0.4),
              size=18, color=RUST, bold=True)
    add_bullets(s, [
        ("AUC 0.681 (95% CI 0.609–0.753) — equivalent to baseline (DeLong p = 0.81)", None),
        ("Brier 0.069 vs 0.228 — a 3.3-fold improvement in calibration", None),
        ("Parametric, interpretable, valid coefficient CIs — regulator-friendly", None),
        ("Net benefit positive in the 5–15% threshold band — the AED-vs-cEEG decision range", None),
    ], Inches(9.0), Inches(2.3), Inches(4.0), Inches(4.6),
       size=12, bullet_size=12)
    add_callout(s,
                 "Same discrimination; far better calibration. The deployment model is Firth.",
                 Inches(0.5), Inches(6.50), Inches(12.3), Inches(0.75),
                 accent=FOREST)
    add_footer(s, prs, slide_num=8)

    # ─────── 9. Conformal — bedside decision support (results 3/5) ───────
    s = blank_slide(prs)
    add_header(s, prs, "Results · 3/5",
                "Conformal prediction turns probabilities into decisions")
    add_image(s, FIG / "F4_conformal.png",
                Inches(0.5), Inches(1.6), w=Inches(8.2))
    add_text(s, "At α = 0.10 (90% coverage)",
              Inches(9.0), Inches(1.7), Inches(4.1), Inches(0.4),
              size=18, color=FOREST, bold=True)
    add_bullets(s, [
        ("Confident singleton prediction in 37% of patients", None),
        ("Rule-out of seizure in 27% — AED prophylaxis can be safely omitted", None),
        ("Rule-in of seizure in 11% — targeted cEEG monitoring", None),
        ("Remaining 63% deferred to clinical judgment (no over-claim)", None),
    ], Inches(9.0), Inches(2.3), Inches(4.0), Inches(4.6),
       size=12, bullet_size=12)
    add_callout(s,
                 "First application of class-conditional conformal sets to postoperative-seizure decision support.",
                 Inches(0.5), Inches(6.50), Inches(12.3), Inches(0.75),
                 accent=FOREST)
    add_footer(s, prs, slide_num=9)

    # ─────── 10. CEA dominance (results 4/5) ───────
    s = blank_slide(prs)
    add_header(s, prs, "Results · 4/5", "ML-guided AED dominates current practice")
    add_image(s, FIG / "F5_cea.png",
                Inches(0.5), Inches(1.6), w=Inches(8.2))
    add_text(s, "Refreshed CEA inputs",
              Inches(9.0), Inches(1.7), Inches(4.1), Inches(0.4),
              size=18, color=NAVY, bold=True)
    add_bullets(s, [
        ("WTP $100k/QALY (Neumann 2014, Vanness 2021, Crespo 2023)", None),
        ("Ceribell-era cEEG cost (Parvizi 2021)", None),
        ("Geriatric AED adverse-event burden (Tsai 2024, Bresser 2022)", None),
        ("ML-AED: $4,365 / 7.43 QALYs — dominant", None),
        ("ML-cEEG: cost-effective in 62% of PSA samples at $100k", None),
    ], Inches(9.0), Inches(2.3), Inches(4.0), Inches(4.6),
       size=12, bullet_size=12)
    add_footer(s, prs, slide_num=10)

    # ─────── 11. VOI — research priorities (results 5/5) ───────
    s = blank_slide(prs)
    add_header(s, prs, "Results · 5/5",
                "Value-of-information identifies the research frontier")
    add_image(s, FIG / "F6_voi.png",
                Inches(0.5), Inches(1.6), w=Inches(8.2))
    add_text(s, "Population EVPI ≈ $190M",
              Inches(9.0), Inches(1.7), Inches(4.1), Inches(0.4),
              size=18, color=OCHRE, bold=True)
    add_bullets(s, [
        ("First VOI applied to postoperative-seizure prophylaxis", None),
        ("Per-patient EVPI at $100k WTP ≈ $541", None),
        ("Top EVPPI parameters (research-priority frontier):", None),
        ("→ cEEG cost-per-day  ·  seizure prevalence  ·  AED RRR  ·  ML sensitivity", None),
    ], Inches(9.0), Inches(2.3), Inches(4.0), Inches(4.6),
       size=12, bullet_size=12)
    add_callout(s,
                 "VOI ranks where to invest in future evidence collection — by decision-relevant value, not by p-value.",
                 Inches(0.5), Inches(6.50), Inches(12.3), Inches(0.75),
                 accent=OCHRE)
    add_footer(s, prs, slide_num=11)

    # ─────── 12. Methodological contributions ───────
    s = blank_slide(prs)
    add_header(s, prs, "Discussion · 1/2", "Methodological contributions")
    add_bullets(s, [
        ("Firth penalized LR as the deployment model for small clinical cohorts",
         "Calibration and interpretability replace AUC as the optimisation target."),
        ("Class-conditional conformal sets for bedside risk stratification",
         "Distribution-free coverage; 27% rule-out at 90% confidence."),
        ("Corrected ICD-10 outcome for nationwide cSDH-seizure analyses",
         "Acute symptomatic seizure ≠ pre-existing epilepsy; the previously-reported population signal disappears under the corrected definition."),
        ("Documented biological transfer-learning failure",
         "eICU age coefficient is negative; pure post-craniotomy cSDH shows a positive age effect. Any future transfer between these populations must verify coefficient signs."),
        ("First value-of-information analysis in postoperative-seizure prevention",
         "Population EVPI $190M / 10yr; per-parameter EVPPI ranks research priorities by decision-relevant value."),
    ], Inches(0.6), Inches(1.6), Inches(12.1), Inches(5.5), size=14)
    add_footer(s, prs, slide_num=12)

    # ─────── 13. Limitations ───────
    s = blank_slide(prs)
    add_header(s, prs, "Discussion · 2/2",
                "Limitations — paired with mitigations")
    add_bullets(s, [
        ("AUC ceiling at 48 events",
         "Bernoulli noise alone gives 95% CI half-width ≈ 0.06. Calibration and clinical utility, not AUC, are the right optimisation target at this sample size."),
        ("Single-institution development cohort",
         "External evaluation in eICU across 139 hospitals with I² = 0% provides the heterogeneity check."),
        ("Imaging features absent from structured EMR",
         "The deployment model uses 21 structured variables; imaging-NLP augmentation is a planned next step."),
        ("Administrative outcome ascertainment",
         "Sensitivity analyses across four time-window cuts (0–24h, 0–48h, 0–72h, ≥24h) preserve the primary AUC."),
        ("US-payer-perspective cost inputs",
         "VOI explicitly identifies cEEG cost as the highest-EVPPI parameter for international refinement."),
    ], Inches(0.6), Inches(1.6), Inches(12.1), Inches(5.5), size=13)
    add_footer(s, prs, slide_num=13)

    # ─────── 14. Conclusions + QR (final) ───────
    s = blank_slide(prs, bg=NAVY)
    # Title block
    add_text(s, "TAKE-HOMES",
              Inches(0.7), Inches(0.6), Inches(8), Inches(0.5),
              size=14, color=RUST, bold=True)
    add_text(s, "Five conclusions in one minute",
              Inches(0.7), Inches(1.0), Inches(8), Inches(0.8),
              size=30, color=WHITE, bold=True)
    add_bullets(s, [
        ("The discrimination ceiling at AUC ≈ 0.68 is biological, not modelling failure.", None),
        ("Calibration replaces AUC as the actionable target — Firth LR matches discrimination, 3× better Brier.", None),
        ("Conformal sets translate the model into bedside decisions for ~37% of patients with 90% coverage.", None),
        ("ML-guided AED is the dominant cost-effectiveness strategy; ML-cEEG is cost-effective at $100k/QALY.", None),
        ("VOI prioritises cEEG cost, seizure prevalence and AED efficacy as the future research frontier.", None),
    ], Inches(0.7), Inches(2.1), Inches(8.0), Inches(4.6),
       size=15, color=WHITE, bullet_size=15, line_spacing=1.5)
    # QR + invitation card
    add_box(s, Inches(9.2), Inches(2.1), Inches(3.4), Inches(4.8),
              fill=WHITE, border=RUST, border_w=2)
    add_text(s, "Read the code · run the calculator",
              Inches(9.4), Inches(2.3), Inches(3.0), Inches(0.5),
              size=14, color=NAVY_DARK, bold=True, align=PP_ALIGN.CENTER)
    add_qr(s, Inches(9.55), Inches(2.85), Inches(2.7))
    add_text(s, "github.com/nielspac177/\ncsdh-postop-seizure-risk",
              Inches(9.3), Inches(5.7), Inches(3.2), Inches(0.7),
              size=11, color=GREY, italic=True,
              align=PP_ALIGN.CENTER, line_spacing=1.3)
    # Thank-you + contact
    add_text(s, "Thank you  ·  nielspacheco1997@gmail.com",
              Inches(0.7), prs.slide_height - Inches(0.7),
              prs.slide_width - Inches(1.4), Inches(0.4),
              size=12, color=NAVY_SOFT, italic=True, align=PP_ALIGN.CENTER)

    prs.save(OUT_PATH)
    print(f"[OK] {OUT_PATH}")
    print(f"     size: {os.path.getsize(OUT_PATH)/1024:.1f} KB")
    print(f"     slides: {len(prs.slides)}  ·  target 14 (15-min talk @ ~1 min/slide + Q&A)")


if __name__ == "__main__":
    build()
