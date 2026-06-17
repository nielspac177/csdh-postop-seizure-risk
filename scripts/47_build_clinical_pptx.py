"""Task 47, Clinical 15-minute talk (PowerPoint).

Audience: neurosurgery / clinical. Message: a calibrated risk tool that knows when to
abstain (rule-out / rule-in / defer), honest about modest discrimination, with a
decision analysis showing selective allocation matters most where AED is of unproven
benefit. Uses the CURRENT numbers and the slide figures from task 46.

Output: Manuscript_05192026/clinical_slides/clinical_talk.pptx
"""
import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

BASE = Path("/Users/nielspacheco/Desktop/Research/Ogilvy research/"
            "Data Chronic Subdural Haematoma/Manuscript_05192026/clinical_slides")
FIG = BASE / "fig"
PAPER_FIG = Path("/Users/nielspacheco/Desktop/Research/Ogilvy research/"
                 "Data Chronic Subdural Haematoma/revision_analyses/figures")

NAVY = RGBColor(0x1F, 0x3B, 0x57)
RUST = RGBColor(0xB5, 0x48, 0x2A)
FOREST = RGBColor(0x3F, 0x6F, 0x4F)
DARK = RGBColor(0x22, 0x22, 0x22)
GREY = RGBColor(0x60, 0x60, 0x60)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xF2, 0xF1, 0xEC)

prs = Presentation()
prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height
FOOTER = "Postoperative seizure after cSDH, a calibrated, conformal decision tool"
_n = [0]


def _box(slide, l, t, w, h):
    tb = slide.shapes.add_textbox(l, t, w, h); tb.text_frame.word_wrap = True
    return tb.text_frame


def _set(p, text, size, color=DARK, bold=False, align=PP_ALIGN.LEFT, italic=False):
    p.text = text; p.alignment = align
    r = p.runs[0]; r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color; r.font.name = "Calibri"


def _bg(slide, color):
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = color


def _rect(slide, l, t, w, h, color):
    from pptx.enum.shapes import MSO_SHAPE
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()
    s.shadow.inherit = False
    return s


def _chrome(slide, title, kicker=None):
    _bg(slide, WHITE)
    _rect(slide, 0, 0, SW, Inches(0.12), NAVY)
    tf = _box(slide, Inches(0.6), Inches(0.35), SW - Inches(1.2), Inches(1.0))
    if kicker:
        p = tf.paragraphs[0]; _set(p, kicker.upper(), 13, RUST, bold=True)
        p2 = tf.add_paragraph(); _set(p2, title, 30, NAVY, bold=True)
    else:
        _set(tf.paragraphs[0], title, 30, NAVY, bold=True)
    # footer
    ff = _box(slide, Inches(0.6), SH - Inches(0.45), SW - Inches(2.0), Inches(0.35))
    _set(ff.paragraphs[0], FOOTER, 9, GREY)
    _n[0] += 1
    nf = _box(slide, SW - Inches(1.0), SH - Inches(0.45), Inches(0.6), Inches(0.35))
    _set(nf.paragraphs[0], str(_n[0]), 11, GREY, align=PP_ALIGN.RIGHT)


def bullets(slide, items, left=Inches(0.7), top=Inches(1.7), width=None, size=20, gap=10):
    width = width or (SW - Inches(1.4))
    tf = _box(slide, left, top, width, SH - top - Inches(0.7))
    for i, (txt, lvl, *style) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        color = style[0] if style else DARK
        bold = style[1] if len(style) > 1 else False
        bullet = "•  " if lvl == 0 else "–  "
        _set(p, ("" if lvl == 0 else "    ") + bullet + txt, size - lvl * 2, color, bold=bold)
        p.space_after = Pt(gap)


def figure_slide(title, img, kicker=None, caption=None, img_w=Inches(8.2)):
    s = prs.slides.add_slide(BLANK); _chrome(s, title, kicker)
    img = Path(img)
    if img.exists():
        pic = s.shapes.add_picture(str(img), 0, Inches(1.75), height=Inches(4.9))
        pic.left = int((SW - pic.width) / 2)
    if caption:
        cf = _box(s, Inches(0.7), SH - Inches(0.95), SW - Inches(1.4), Inches(0.5))
        _set(cf.paragraphs[0], caption, 13, GREY, italic=True, align=PP_ALIGN.CENTER)
    return s


def split_slide(title, bullet_items, img, kicker=None, caption=None):
    s = prs.slides.add_slide(BLANK); _chrome(s, title, kicker)
    bullets(s, bullet_items, left=Inches(0.7), top=Inches(1.85),
            width=Inches(5.7), size=18, gap=8)
    img = Path(img)
    if img.exists():
        pic = s.shapes.add_picture(str(img), Inches(6.7), Inches(1.85), width=Inches(6.2))
    if caption:
        cf = _box(s, Inches(6.7), SH - Inches(0.95), Inches(6.2), Inches(0.5))
        _set(cf.paragraphs[0], caption, 12, GREY, italic=True, align=PP_ALIGN.CENTER)
    return s


def divider(title, sub):
    s = prs.slides.add_slide(BLANK); _bg(s, NAVY)
    _rect(s, Inches(4.67), Inches(3.0), Inches(4.0), Emu(38100), RUST)
    tf = _box(s, Inches(1.0), Inches(3.2), SW - Inches(2.0), Inches(1.6))
    _set(tf.paragraphs[0], title, 40, WHITE, bold=True, align=PP_ALIGN.CENTER)
    p = tf.add_paragraph(); _set(p, sub, 18, RGBColor(0xAF, 0xC0, 0xCF), align=PP_ALIGN.CENTER)
    _n[0] += 1
    return s


# ── 1. Title ──
s = prs.slides.add_slide(BLANK); _bg(s, NAVY)
_rect(s, 0, Inches(2.5), SW, Inches(0.04), RUST)
tf = _box(s, Inches(0.9), Inches(2.7), SW - Inches(1.8), Inches(2.2))
_set(tf.paragraphs[0], "Predicting seizures after chronic subdural haematoma surgery",
     34, WHITE, bold=True)
p = tf.add_paragraph(); _set(p, "A calibrated risk tool that knows when to abstain, "
                             "and what it means for AED prophylaxis", 20, RGBColor(0xCF, 0xDA, 0xE5))
tf2 = _box(s, Inches(0.9), Inches(5.4), SW - Inches(1.8), Inches(1.3))
_set(tf2.paragraphs[0], "Niels Pacheco-Barrios, MD", 18, WHITE, bold=True)
p = tf2.add_paragraph(); _set(p, "Department of Neurosurgery, Beth Israel Deaconess Medical Center · "
                             "Harvard Medical School", 14, RGBColor(0xCF, 0xDA, 0xE5))
p = tf2.add_paragraph(); _set(p, "15-minute presentation", 12, RGBColor(0x9F, 0xB3, 0xC4), italic=True)

divider("Introduction", "The clinical problem and why current practice is uncertain")

# ── 2. The clinical problem ──
s = prs.slides.add_slide(BLANK); _chrome(s, "The clinical dilemma", "Background")
bullets(s, [
    ("Postoperative seizures complicate 7–12% of cSDH evacuations", 0, DARK, True),
    ("Associated with longer ICU stay, higher 30-day mortality, worse function", 0),
    ("So should we give antiepileptic drugs (AED) to everyone after surgery?", 0, NAVY, True),
    ("AEDs in the elderly are not benign:", 0),
    ("falls (risk ↑ up to 2×), sedation, cognitive slowing, drug interactions", 1, RUST),
    ("Treating everyone exposes the ~90% who never seize to that harm", 0),
    ("We lack a way to say, for THIS patient, how much seizure risk there is", 0, NAVY, True),
])

# ── 3. AED evidence ──
figure_slide("Routine AED prophylaxis after cSDH is unproven", FIG / "sld_aed_evidence.png",
             kicker="Why this matters",
             caption="No randomised trial; pooled and adjusted estimates show no significant seizure reduction.")

divider("Methods", "Cohorts, model, and decision analysis")

# ── 4. What we built ──
s = prs.slides.add_slide(BLANK); _chrome(s, "What we built", "Approach")
bullets(s, [
    ("A risk model optimised for CALIBRATION and DECISION SUPPORT, not for chasing AUC", 0, NAVY, True),
    ("Firth penalised logistic regression on 18 variables available at the end of surgery", 0),
    ("(before the AED/EEG decision, no information leakage)", 1, GREY),
    ("A conformal layer that returns rule-out / rule-in / or 'defer to clinician'", 0, FOREST, True),
    ("A cost-effectiveness + value-of-information analysis linking risk to action", 0),
])

# ── 4b. Design & cohorts ──
s = prs.slides.add_slide(BLANK); _chrome(s, "Design and cohorts", "Methods")
bullets(s, [
    ("Development: BIDMC, 655 cSDH evacuations, 48 postoperative seizures (2010–2023)", 0, NAVY, True),
    ("External evaluation: eICU, 3,297 non-traumatic subdural ICU stays across 42 hospitals (300 seizures), from a 5,376-stay screen", 0, NAVY, True),
    ("5×5 repeated cross-validation; Platt calibration; class-conditional conformal sets", 0),
    ("Decision tree + 10-year Markov; probabilistic sensitivity (10,000 iterations)", 0),
    ("eICU is a related, mixed-acuity ICU population, not operative cSDH — a deliberate transportability test", 0, GREY),
])

divider("Results", "Model performance, abstention, and the decision analysis")

# ── 5. The model & the honest ceiling ──
s = figure_slide("Discrimination is modest, and that is the point", FIG / "sld_ceiling.png",
             kicker="The model", img_w=Inches(9.0),
             caption="Flexible models 'win' AUC only by over-fitting 48 events; we optimise calibration "
                     "and honest uncertainty, what a bedside tool actually needs.")

# ── 6. Calibration ──
figure_slide("Calibrated on average, deliberately conservative", FIG / "calibration.png",
             kicker="Calibration",
             caption="When the model says '5%', about 5% of such patients seize; predictions are "
                     "deliberately compressed (conservative) given only 48 events.")

# ── 7. Conformal, knows when to abstain ──
figure_slide("A tool that knows when it doesn't know", FIG / "sld_conformal.png",
             kicker="Conformal prediction",
             caption="Conformal = a guarantee that ~90% of 'rule-out' patients truly won't seize; "
                     "the tool gives an actionable answer in ~22% and openly defers the rest.")

# ── 7b. Bedside example ──
s = prs.slides.add_slide(BLANK); _chrome(s, "What it looks like at the bedside", "Worked example (illustrative)")
bullets(s, [
    ("At end of evacuation, enter routine variables:", 0, NAVY, True),
    ("Patient A, older, MMA-embolisation, decompression, no prior seizure", 0, FOREST, True),
    ("→ low risk → rule-OUT → observe, spare AED", 1, FOREST),
    ("Patient B, re-accumulation needing drainage, dense collection", 0, RUST, True),
    ("→ high risk → rule-IN → cEEG + targeted AED", 1, RUST),
    ("Most patients fall in between → the tool defers to judgement", 0, GREY),
], left=Inches(0.7), top=Inches(1.85), width=Inches(6.0), size=17, gap=8)
bullets(s, [
    ("What moves the estimate", 0, NAVY, True),
    ("↑ drainage / re-accumulation, denser collection, midline shift", 1, DARK),
    ("↓ MMA embolisation, surgical decompression", 1, DARK),
    ("Strong predictors are procedure variables, they reflect which "
     "operation was needed (confounding by indication), so this is a "
     "decision-support prompt, not an autonomous rule.", 0, GREY),
], left=Inches(7.0), top=Inches(1.85), width=Inches(5.8), size=15, gap=10)

# ── 8. From prediction to decision ──
split_slide("From risk to decision: does the model change care?",
    [("Four strategies: observe / AED-for-all / ML-guided AED / ML-guided cEEG", 0, DARK, True),
     ("All depend on one unknown: does AED actually prevent cSDH seizures?", 0, NAVY, True),
     ("If AED works well → treat everyone", 0, RUST),
     ("If AED barely works (the cSDH reality) → use the model to spare low-risk patients", 0, FOREST, True),
     ("Treating everyone can be actively harmful:", 0, RUST, True),
     ("when AED is ineffective, universal AED is WORSE than simple observation", 1, RUST),
     ("ML-guided allocation beats observation at every efficacy value", 0, FOREST)],
    FIG / "sld_cea_curve.png", kicker="Cost-effectiveness",
    caption="Crossover near RRR 0.30; cSDH-plausible range favours ML-guided allocation.")

# ── 9. Model adds value beyond treating fewer ──
figure_slide("It's the model, not just treating fewer people", FIG / "sld_premium.png",
             kicker="Does discrimination matter?",
             caption="Vs random allocation at the same treated fraction, the model keeps a positive net-benefit premium.")

# ── 10. Value of information ──
s = prs.slides.add_slide(BLANK); _chrome(s, "What should we study next?", "Value of information")
bullets(s, [
    ("The decision hinges on two numbers we don't yet have for cSDH:", 0, NAVY, True),
    ("how well AED prevents seizures (efficacy)", 1, RUST),
    ("how much AED harms elderly patients (disutility)", 1, RUST),
    ("Value of information, the expected payoff of running the trial that answers this, ≈ $23M over 10 years (US)", 0, DARK, True),
    ("The model's discrimination is honestly a smaller lever than the AED question itself", 0, GREY),
    ("Priority: a focused cSDH AED-prophylaxis trial + per-day cEEG cost data", 0, FOREST, True),
])

divider("Discussion & limitations", "What this is, what it isn't, and what comes next")

# ── 11. Limitations ──
s = prs.slides.add_slide(BLANK); _chrome(s, "What this is, and isn't", "Honest limitations")
bullets(s, [
    ("Single-centre development (48 events), a proof of concept, not a finished tool", 0, RUST, True),
    ("External cohort (eICU) is ICU subdural patients, not all operative cSDH", 0),
    ("Discrimination is modest; predictions are conservative (under-dispersed)", 0),
    ("Outcome is chart/structured-flag, not EEG-adjudicated", 0),
    ("What IS solid: honest calibration, abstention, and a transparent decision framework", 0, FOREST, True),
    ("Next: prospective, multi-centre validation with site recalibration", 0, NAVY, True),
])

# ── 12. Take-home ──
s = prs.slides.add_slide(BLANK); _bg(s, NAVY)
_rect(s, 0, Inches(0.0), SW, Inches(0.12), RUST)
tf = _box(s, Inches(0.9), Inches(0.9), SW - Inches(1.8), Inches(0.8))
_set(tf.paragraphs[0], "Take-home", 30, WHITE, bold=True)
tf2 = _box(s, Inches(0.9), Inches(2.0), SW - Inches(1.8), Inches(4.6))
takeaways = [
    "A calibrated risk model can give trustworthy, individual seizure probabilities after cSDH surgery.",
    "Conformal prediction lets the tool abstain, confident rule-out/rule-in for ~22%, defer for the rest.",
    "Because cSDH AED efficacy is unproven, treating everyone can do net harm, universal AED was worse "
    "than observation when AED is ineffective. Selective, risk-guided use is the rational target.",
    "The decision is governed by AED efficacy and harm, the priorities for the next trial.",
]
for i, t in enumerate(takeaways):
    p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
    _set(p, "→  " + t, 19, WHITE, bold=(i in (1, 2))); p.space_after = Pt(16)
tf3 = _box(s, Inches(0.9), SH - Inches(1.0), SW - Inches(1.8), Inches(0.6))
_set(tf3.paragraphs[0], "Code, figures, and reproducibility: github.com/nielspac177/csdh-postop-seizure-risk",
     13, RGBColor(0x9F, 0xB3, 0xC4), italic=True)

out = BASE / "clinical_talk.pptx"
prs.save(out)
print(f"[OK] {out}  ({os.path.getsize(out)/1024:.0f} KB, {len(prs.slides._sldIdLst)} slides)")
